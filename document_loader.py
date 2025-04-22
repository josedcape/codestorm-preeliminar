"""
Módulo para cargar y procesar documentos de diferentes formatos.
Permite extraer texto de distintos tipos de archivos para proporcionar contexto a los agentes.
"""

import os
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Union

# Importar librerías para procesamiento de diferentes tipos de documentos
from bs4 import BeautifulSoup
from markdown import markdown
import pypdf
import docx
from pptx import Presentation
import ebooklib
from ebooklib import epub

logger = logging.getLogger(__name__)

class DocumentLoader:
    """Clase para cargar y procesar documentos de diferentes formatos."""
    
    @staticmethod
    def load_document(file_path: str) -> Optional[str]:
        """
        Carga y extrae el contenido textual de un documento.
        
        Args:
            file_path: Ruta al archivo
            
        Returns:
            str: Contenido del documento o None si no se pudo procesar
        """
        if not os.path.exists(file_path):
            logger.error(f"El archivo no existe: {file_path}")
            return None
            
        try:
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.pdf':
                return DocumentLoader.extract_text_from_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                return DocumentLoader.extract_text_from_docx(file_path)
            elif ext in ['.html', '.htm']:
                return DocumentLoader.extract_text_from_html(file_path)
            elif ext == '.md':
                return DocumentLoader.extract_text_from_markdown(file_path)
            elif ext in ['.txt', '.py', '.js', '.css', '.json', '.csv', '.xml', '.yml', '.yaml']:
                return DocumentLoader.extract_text_from_text(file_path)
            elif ext == '.pptx':
                return DocumentLoader.extract_text_from_pptx(file_path)
            elif ext == '.epub':
                return DocumentLoader.extract_text_from_epub(file_path)
            else:
                logger.warning(f"Formato de archivo no soportado: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error al procesar el documento {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extrae texto de un archivo PDF."""
        text = ""
        
        try:
            with open(file_path, 'rb') as file:
                reader = pypdf.PdfReader(file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
                    
            # Limpiar espacios en blanco excesivos
            text = re.sub(r'\s+', ' ', text).strip()
            return text
        except Exception as e:
            logger.error(f"Error al extraer texto del PDF {file_path}: {str(e)}")
            return f"Error al procesar el PDF: {str(e)}"
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extrae texto de un archivo DOCX."""
        try:
            doc = docx.Document(file_path)
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            logger.error(f"Error al extraer texto del DOCX {file_path}: {str(e)}")
            return f"Error al procesar el documento DOCX: {str(e)}"
    
    @staticmethod
    def extract_text_from_html(file_path: str) -> str:
        """Extrae texto de un archivo HTML."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                
                # Eliminar scripts y estilos
                for script in soup(["script", "style"]):
                    script.extract()
                
                # Extraer texto
                text = soup.get_text(separator=' ')
                
                # Limpiar espacios en blanco excesivos
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return text
        except Exception as e:
            logger.error(f"Error al extraer texto del HTML {file_path}: {str(e)}")
            return f"Error al procesar el HTML: {str(e)}"
    
    @staticmethod
    def extract_text_from_markdown(file_path: str) -> str:
        """Extrae texto de un archivo Markdown."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_text = file.read()
                
                # Convertir a HTML y luego extraer el texto
                html = markdown(md_text)
                soup = BeautifulSoup(html, 'html.parser')
                text = soup.get_text()
                
                return text
        except Exception as e:
            logger.error(f"Error al extraer texto del Markdown {file_path}: {str(e)}")
            return f"Error al procesar el Markdown: {str(e)}"
    
    @staticmethod
    def extract_text_from_text(file_path: str) -> str:
        """Extrae texto de archivos de texto plano."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error al extraer texto del archivo {file_path}: {str(e)}")
            return f"Error al procesar el archivo de texto: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extrae texto de una presentación PowerPoint."""
        try:
            presentation = Presentation(file_path)
            text = ""
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
                
            return text
        except Exception as e:
            logger.error(f"Error al extraer texto del PPTX {file_path}: {str(e)}")
            return f"Error al procesar la presentación PowerPoint: {str(e)}"
    
    @staticmethod
    def extract_text_from_epub(file_path: str) -> str:
        """Extrae texto de un libro electrónico EPUB."""
        try:
            book = epub.read_epub(file_path)
            text = ""
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    html_content = item.get_content().decode('utf-8')
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text += soup.get_text() + "\n\n"
            
            return text
        except Exception as e:
            logger.error(f"Error al extraer texto del EPUB {file_path}: {str(e)}")
            return f"Error al procesar el libro electrónico: {str(e)}"


def get_document_summary(file_path: str, max_length: int = 1000) -> Dict:
    """
    Obtiene un resumen de un documento para mostrar en la interfaz.
    
    Args:
        file_path: Ruta al archivo
        max_length: Longitud máxima del texto extraído para el resumen
        
    Returns:
        dict: Información sobre el documento incluyendo tipo y extracto
    """
    try:
        if not os.path.exists(file_path):
            return {"success": False, "error": "El archivo no existe"}
            
        file_size = os.path.getsize(file_path)
        file_type = os.path.splitext(file_path)[1].lower()
        file_name = os.path.basename(file_path)
        
        text = DocumentLoader.load_document(file_path)
        
        if text:
            # Truncar el texto para el resumen
            text_preview = text[:max_length] + "..." if len(text) > max_length else text
            total_words = len(text.split())
            
            return {
                "success": True,
                "file_name": file_name,
                "file_type": file_type,
                "file_size": file_size,
                "total_words": total_words,
                "text_preview": text_preview
            }
        else:
            return {"success": False, "error": "No se pudo extraer texto del documento"}
    except Exception as e:
        logger.error(f"Error al obtener resumen del documento {file_path}: {str(e)}")
        return {"success": False, "error": str(e)}


def load_document_as_context(file_path: str, max_length: int = 10000) -> Dict:
    """
    Carga un documento y lo prepara como contexto para los agentes.
    
    Args:
        file_path: Ruta al archivo
        max_length: Longitud máxima del texto para usar como contexto
        
    Returns:
        dict: Contexto formateado para los agentes
    """
    try:
        text = DocumentLoader.load_document(file_path)
        
        if not text:
            return {
                "success": False,
                "error": "No se pudo extraer texto del documento"
            }
        
        # Truncar si es demasiado largo
        if len(text) > max_length:
            text = text[:max_length] + "... [Contenido truncado debido a su longitud]"
        
        file_name = os.path.basename(file_path)
        file_type = os.path.splitext(file_path)[1].lower()
        
        return {
            "success": True,
            "context": {
                "source": file_name,
                "type": file_type,
                "content": text,
                "word_count": len(text.split())
            }
        }
    except Exception as e:
        logger.error(f"Error al cargar documento como contexto {file_path}: {str(e)}")
        return {"success": False, "error": str(e)}